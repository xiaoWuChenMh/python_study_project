# pip install python-pptx
from pptx import Presentation
from PIL import Image
from io import BytesIO
from tool.image.image_util import divide_image_ratio,save_divided_images
class PptReadImage:

    def __init__(self):
        print("read ppt")

    def extract_images_from_ppt(self,ppt_path):
        image_index = 0
        presentation = Presentation(ppt_path)
        for slide_number, slide in enumerate(presentation.slides, start=1):
            print(f"Processing Slide {slide_number}")
            for shape in slide.shapes:
                # 提取图片
                image = shape.image
                # 使用Pillow打开图片
                image_bytes = BytesIO(image.blob)
                width, height = image.size
                img = Image.open(image_bytes).crop((0,4300,width+2,height))
                # img.show()  # 预览图片


                divided_images = divide_image_ratio(img, 4,10)
                # 保存切割后的图片
                base_name = f"地址\slide_{slide_number}_image_{image_index}"
                save_divided_images(divided_images, base_name, image_index)
                image_index += 1


if __name__ == '__main__':
    read_image = PptReadImage()
    # 使用函数
    ppt_path = "地址\Midjourney全套颜色关键词.pptx"  # 你的PDF文件路径
    read_image.extract_images_from_ppt(ppt_path)
